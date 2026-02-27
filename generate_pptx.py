#!/usr/bin/env python3
"""
Gera apresentação PowerPoint para o projeto HashCrackerLab.
Projeto Final — Cibersegurança (Fevereiro 2026)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Cores do tema ──
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # fundo escuro
BG_SLIDE     = RGBColor(0x10, 0x17, 0x20)   # fundo slide
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)   # ciano principal
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)   # verde sucesso
ACCENT_RED   = RGBColor(0xFF, 0x45, 0x45)   # vermelho alerta
ACCENT_ORANGE= RGBColor(0xFF, 0xA5, 0x00)   # laranja
ACCENT_PURPLE= RGBColor(0xBB, 0x86, 0xFC)   # roxo
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB8, 0xC0)
DARK_GRAY    = RGBColor(0x1A, 0x23, 0x2E)
MID_GRAY     = RGBColor(0x2A, 0x35, 0x42)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Define cor de fundo do slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    """Adiciona retângulo colorido."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val=str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Adiciona caixa de texto."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_CYAN, font_name="Segoe UI"):
    """Adiciona lista com bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        # bullet char
        run_b = p.add_run()
        run_b.text = "▸ "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name
        # text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = font_name
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=ACCENT_CYAN, cell_color=LIGHT_GRAY):
    """Adiciona tabela estilizada."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            # Estilo
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Segoe UI"
                    if row_idx == 0:
                        run.font.size = Pt(13)
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.size = Pt(12)
                        run.font.color.rgb = cell_color
            # Background das células
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = RGBColor(0x15, 0x1F, 0x2B)
            elif row_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x24)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x0E, 0x15, 0x1E)
    return table_shape


def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    """Linha decorativa horizontal."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, accent=ACCENT_CYAN):
    """Card com borda de acento."""
    # Background
    bg = add_shape_fill(slide, left, top, width, height, DARK_GRAY)
    # Accent bar left
    add_shape_fill(slide, left, top, Pt(4), height, accent)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.1),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=15, color=accent, bold=True)
    # Body
    y = top + Inches(0.5)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.3), y,
                     width - Inches(0.4), Inches(0.3),
                     line, font_size=12, color=LIGHT_GRAY)
        y += Inches(0.25)


# ═══════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Capa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent line top
    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_CYAN)

    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
                 "HashCrackerLab", font_size=52, color=WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
                 "Plataforma Integrada de Cibersegurança Ofensiva",
                 font_size=24, color=ACCENT_CYAN)

    add_accent_line(slide, Inches(1), Inches(3.7), Inches(3), ACCENT_CYAN)

    # Description
    add_text_box(slide, Inches(1), Inches(4.1), Inches(10), Inches(0.6),
                 "Cracking de Hashes · WiFi WPA2 · Telnet Credential Harvesting",
                 font_size=18, color=LIGHT_GRAY)

    # Team
    add_text_box(slide, Inches(1), Inches(5.2), Inches(10), Inches(0.5),
                 "Henrique Carvalho  ·  Gonçalo Ferro  ·  Francisco Silva  ·  Duarte Vilar",
                 font_size=16, color=LIGHT_GRAY)

    # Context
    add_text_box(slide, Inches(1), Inches(5.8), Inches(10), Inches(0.4),
                 "Projeto Final — Cibersegurança  |  Fevereiro 2026",
                 font_size=14, color=RGBColor(0x70, 0x78, 0x80))

    # Bottom bar
    add_shape_fill(slide, Inches(0), Inches(7.2), SLIDE_W, Pt(4), ACCENT_CYAN)


def slide_agenda(prs):
    """Slide 2 — Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Agenda", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2))

    items = [
        ("01", "Contexto e Motivação", "Porquê este projeto?"),
        ("02", "Arquitetura do Laboratório", "Rede isolada, 4 máquinas, 3 vetores"),
        ("03", "Demo: WiFi WPA2 Cracking", "Captura de handshake + cracking"),
        ("04", "Demo: Telnet Credential Capture", "Wireshark — credenciais em texto claro"),
        ("05", "Demo: Hash Cracking CPU vs GPU", "15 passwords × 4 algoritmos × 5 ataques"),
        ("06", "Resultados e Análise", "Tabela comparativa + speedup GPU"),
        ("07", "Conclusões", "Algoritmo forte + password forte = segurança"),
    ]

    y = Inches(1.7)
    for num, title, desc in items:
        # Number circle
        add_text_box(slide, Inches(1), y, Inches(0.6), Inches(0.5),
                     num, font_size=20, color=ACCENT_CYAN, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.7), y - Inches(0.02), Inches(4), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.7), y + Inches(0.32), Inches(6), Inches(0.3),
                     desc, font_size=13, color=LIGHT_GRAY)
        y += Inches(0.72)


def slide_context(prs):
    """Slide 3 — Contexto e Motivação"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Contexto e Motivação", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    # Pergunta de investigação
    add_shape_fill(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1), DARK_GRAY)
    add_text_box(slide, Inches(1.1), Inches(1.7), Inches(11), Inches(0.8),
                 '❝ Quão eficaz é a aceleração por GPU no cracking de passwords e como\n'
                 '  a escolha do algoritmo de hashing influencia a resistência a ataques? ❞',
                 font_size=17, color=ACCENT_CYAN, bold=False)

    # Problema
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.4),
                 "O Problema", font_size=22, color=ACCENT_RED, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(2), [
        "Demos de segurança são manuais e não reprodutíveis",
        "Sem forma standard de comparar algoritmos de hashing",
        "Relação algoritmo ↔ hardware ↔ tempo raramente quantificada",
        "Ferramentas isoladas sem contexto nem métricas",
    ], bullet_color=ACCENT_RED)

    # Solução
    add_text_box(slide, Inches(6.8), Inches(3.0), Inches(6), Inches(0.4),
                 "A Nossa Solução", font_size=22, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(2), [
        "Orquestrador Python — pipeline 100% automatizado",
        "Reprodutibilidade: seed fixa + salts determinísticos",
        "Multi-vetor: WiFi + Telnet + Hash Cracking",
        "Relatório comparativo CPU vs GPU automático",
    ], bullet_color=ACCENT_GREEN)


def slide_architecture(prs):
    """Slide 4 — Arquitetura de Rede"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Arquitetura do Laboratório", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
                 "Rede Isolada — 192.168.100.0/24  |  Sem acesso à Internet  |  4 Máquinas",
                 font_size=15, color=LIGHT_GRAY)

    # Router - center top
    add_card(slide, Inches(4.5), Inches(2.2), Inches(4.2), Inches(1.0),
             "Router TP-Link  —  192.168.100.1",
             ["SSID: LAB-SERVERS  |  WPA2-PSK AES"], accent=ACCENT_ORANGE)

    # Henrique - left
    add_card(slide, Inches(0.5), Inches(3.8), Inches(3.5), Inches(1.5),
             "Henrique — Arch Linux (.10)",
             ["Orchestrator + GPU Cracking", "NVIDIA OpenCL", "Ethernet"], accent=ACCENT_CYAN)

    # Ferro - right
    add_card(slide, Inches(9.2), Inches(3.8), Inches(3.5), Inches(1.5),
             "Ferro — Kali Linux (.20)",
             ["WiFi Monitor Mode", "Aircrack-ng + Hashcat", "Captura WPA2"], accent=ACCENT_GREEN)

    # Francisco - bottom left
    add_card(slide, Inches(2), Inches(5.8), Inches(3.8), Inches(1.3),
             "Francisco — Windows (.30)",
             ["Telnet Server + Wireshark", "Captura de pacotes"], accent=ACCENT_PURPLE)

    # Duarte - bottom right
    add_card(slide, Inches(7.3), Inches(5.8), Inches(3.8), Inches(1.3),
             "Duarte — Windows (.40)",
             ["Telnet Client", "Geração de tráfego"], accent=ACCENT_PURPLE)


def slide_pipeline(prs):
    """Slide 5 — Pipeline de Execução"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Pipeline de Execução", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    # Pipeline steps
    steps = [
        ("Config YAML", "ACCENT_CYAN", "Definição de\nparâmetros"),
        ("Orchestrator", "ACCENT_CYAN", "Coordena\ntodo o pipeline"),
        ("Hash Generator", "ACCENT_GREEN", "15 pw × 4 algos\n= 60 hashes"),
        ("Cracking\nGPU + CPU", "ACCENT_ORANGE", "5 modos\nde ataque"),
        ("Metrics", "ACCENT_PURPLE", "CSV + JSON\nagregação"),
        ("Report", "ACCENT_GREEN", "REPORT.md\n+ gráficos"),
    ]

    colors_map = {
        "ACCENT_CYAN": ACCENT_CYAN,
        "ACCENT_GREEN": ACCENT_GREEN,
        "ACCENT_ORANGE": ACCENT_ORANGE,
        "ACCENT_PURPLE": ACCENT_PURPLE,
    }

    x_start = Inches(0.5)
    y_mid = Inches(3.0)
    box_w = Inches(1.8)
    box_h = Inches(1.3)
    gap = Inches(0.3)

    for i, (label, col_name, desc) in enumerate(steps):
        x = x_start + i * (box_w + gap)
        col = colors_map[col_name]
        # Box
        shape = add_shape_fill(slide, x, y_mid, box_w, box_h, DARK_GRAY)
        # Top accent
        add_shape_fill(slide, x, y_mid, box_w, Pt(3), col)
        # Label
        add_text_box(slide, x, y_mid + Inches(0.15), box_w, Inches(0.5),
                     label, font_size=14, color=col, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Desc
        add_text_box(slide, x, y_mid + Inches(0.6), box_w, Inches(0.6),
                     desc, font_size=11, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
        # Arrow
        if i < len(steps) - 1:
            add_text_box(slide, x + box_w, y_mid + Inches(0.35), gap + Inches(0.1), Inches(0.4),
                         "→", font_size=22, color=ACCENT_CYAN,
                         alignment=PP_ALIGN.CENTER)

    # Bottom: componentes
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
                 "Módulos Core", font_size=20, color=WHITE, bold=True)

    modules = [
        ("hash_generator.py", "7 algoritmos (MD5 a Argon2)", ACCENT_GREEN),
        ("cracking_manager.py", "4 modos de ataque hashcat", ACCENT_ORANGE),
        ("metrics_collector.py", "CSV + JSON export", ACCENT_PURPLE),
        ("network_manager.py", "Validação de isolamento", ACCENT_CYAN),
        ("cleanup_manager.py", "3-pass secure delete", ACCENT_RED),
    ]

    x = Inches(0.5)
    for name, desc, col in modules:
        add_shape_fill(slide, x, Inches(5.4), Inches(2.3), Inches(0.8), DARK_GRAY)
        add_shape_fill(slide, x, Inches(5.4), Inches(2.3), Pt(2), col)
        add_text_box(slide, x + Inches(0.1), Inches(5.45), Inches(2.1), Inches(0.3),
                     name, font_size=11, color=col, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(5.75), Inches(2.1), Inches(0.3),
                     desc, font_size=10, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
        x += Inches(2.5)


def slide_algorithms(prs):
    """Slide 6 — Algoritmos de Hashing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Algoritmos de Hashing", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    data = [
        ["Algoritmo", "Hashcat Mode", "Tipo", "Resistência", "Velocidade GPU"],
        ["MD5", "0 / 20", "Hash simples", "Muito Baixa", "22.5 GH/s"],
        ["SHA-256", "1400 / 1420", "Hash + salt", "Baixa", "~4 GH/s"],
        ["Bcrypt", "3200", "Adaptativo (cost)", "Alta", "~26 KH/s"],
        ["Argon2id", "34000", "Memory-hard", "Muito Alta", "~140 H/s"],
    ]

    add_table(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(2.5),
              len(data), len(data[0]), data)

    # Visual comparison - speed bars
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
                 "Velocidade Relativa (GPU)", font_size=18, color=WHITE, bold=True)

    algos = [
        ("MD5", 10, "22.5 GH/s", ACCENT_RED),
        ("SHA-256", 4, "~4 GH/s", ACCENT_ORANGE),
        ("Bcrypt", 0.05, "~26 KH/s", ACCENT_GREEN),
        ("Argon2id", 0.001, "~140 H/s", ACCENT_CYAN),
    ]

    y = Inches(5.0)
    max_bar = Inches(8)
    for name, rel_speed, speed_txt, color in algos:
        add_text_box(slide, Inches(0.8), y, Inches(1.5), Inches(0.35),
                     name, font_size=14, color=WHITE, bold=True)
        bar_width = max(Pt(8), int(max_bar * (rel_speed / 10)))
        add_shape_fill(slide, Inches(2.5), y + Inches(0.05), bar_width, Inches(0.25), color)
        add_text_box(slide, Inches(2.5) + bar_width + Inches(0.15), y, Inches(2), Inches(0.35),
                     speed_txt, font_size=12, color=LIGHT_GRAY)
        y += Inches(0.45)


def slide_attack_modes(prs):
    """Slide 7 — Modos de Ataque"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "5 Modos de Ataque", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    data = [
        ["Modo", "Flag Hashcat", "Descrição", "Keyspace"],
        ["Dicionário", "-a 0", "rockyou.txt (14.3M passwords)", "14.3M"],
        ["Dict + Rules", "-a 0 -r best66.rule", "Mutações: P→p, a→@, +123", "~950M"],
        ["Brute-force PIN", "-a 3 ?d?d?d?d", "PINs 0000–9999", "10K"],
        ["Brute-force Padrão", "-a 3 ?u?l?l?l?d?d", "Ulll99 pattern", "~11.8M"],
        ["Híbrido", "-a 6 wordlist ?d?d?d", "Wordlist + sufixo 000–999", "~14.3B"],
    ]

    add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(3),
              len(data), len(data[0]), data)

    # Passwords sample
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
                 "Amostra de Passwords (15 total)", font_size=18, color=WHITE, bold=True)

    pw_data = [
        ["Categoria", "Qtd", "Exemplos", "Expectativa"],
        ["Fracas (Top 20)", "5", "123456, password, qwerty", "Crackeadas em todos os algoritmos"],
        ["Médias", "5", "summer2024, hunter42", "Crackeadas em MD5/SHA-256"],
        ["Fortes", "5", "X7k#mP9$vL2@, Cr¥pt0_L4b_99", "Resistem a todos os ataques"],
    ]

    add_table(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
              len(pw_data), len(pw_data[0]), pw_data)


def slide_demo_wifi(prs):
    """Slide 8 — Demo WiFi WPA2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Demo: WiFi WPA2 Cracking", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_GREEN)

    # Steps
    steps_data = [
        ("1", "Monitor Mode", "airmon-ng start wlan00", "Capturar todos os pacotes WiFi"),
        ("2", "Scan de Redes", "wifi_cracker.py --scan-only", "Encontrar LAB-SERVERS + BSSID"),
        ("3", "Captura Handshake", "wifi_cracker.py --capture", "Aguardar 4-way handshake WPA2"),
        ("4", "Deauth Attack", "wifi_cracker.py --deauth", "Forçar reconexão dos clientes"),
        ("5", "Cracking", "wifi_cracker.py --crack", "Testar wordlist contra handshake"),
    ]

    y = Inches(1.8)
    for num, title, cmd, desc in steps_data:
        # Step number
        add_text_box(slide, Inches(0.8), y, Inches(0.4), Inches(0.4),
                     num, font_size=22, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, Inches(1.4), y, Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, Inches(4.2), y, Inches(4.5), Inches(0.35),
                     cmd, font_size=13, color=ACCENT_CYAN, font_name="Consolas")
        add_text_box(slide, Inches(9), y, Inches(4), Inches(0.35),
                     desc, font_size=12, color=LIGHT_GRAY)
        y += Inches(0.55)

    # Result box
    add_shape_fill(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.2), DARK_GRAY)
    add_shape_fill(slide, Inches(0.8), Inches(4.8), Pt(4), Inches(1.2), ACCENT_GREEN)
    add_text_box(slide, Inches(1.2), Inches(4.9), Inches(10), Inches(0.4),
                 "Resultado", font_size=16, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.6),
                 "[+] PASSWORD ENCONTRADA: Cibersegura\n[+] Tempo: 3.2 segundos",
                 font_size=15, color=WHITE, font_name="Consolas")

    # Lesson
    add_shape_fill(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7), RGBColor(0x1A, 0x0A, 0x0A))
    add_text_box(slide, Inches(1.2), Inches(6.35), Inches(10.5), Inches(0.6),
                 "⚠  Password com palavra e maiúscula cai em segundos. Proteção: passwords aleatórias ≥16 chars.",
                 font_size=14, color=ACCENT_ORANGE)


def slide_demo_telnet(prs):
    """Slide 9 — Demo Telnet"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Demo: Telnet Credential Capture", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_PURPLE)

    # Left: What happens
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
                 "O Que Acontece", font_size=20, color=WHITE, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.5), [
        "Francisco inicia servidor Telnet (porta 23)",
        "Francisco abre Wireshark (filtro: tcp.port == 23)",
        "Duarte liga-se como admin / SecurePass123",
        "Wireshark captura tudo em texto claro",
    ], bullet_color=ACCENT_PURPLE, font_size=15)

    # Right: Wireshark output
    add_shape_fill(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(3), DARK_GRAY)
    add_shape_fill(slide, Inches(6.5), Inches(1.8), Inches(6), Pt(2), ACCENT_PURPLE)
    add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(0.3),
                 "Wireshark — Pacotes Capturados", font_size=13, color=ACCENT_PURPLE, bold=True)
    wireshark_output = (
        '#42  .40 → .30  Telnet: "admin"\n'
        '#43  .40 → .30  Telnet: "SecurePass123"\n\n'
        '⚠ Password visível em plaintext!'
    )
    add_text_box(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(2),
                 wireshark_output, font_size=14, color=ACCENT_RED, font_name="Consolas")

    # Bottom: lesson
    add_shape_fill(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), DARK_GRAY)
    add_shape_fill(slide, Inches(0.8), Inches(5.2), Pt(4), Inches(1.5), ACCENT_GREEN)
    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10.5), Inches(0.4),
                 "Solução", font_size=18, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(1.2), Inches(5.8), Inches(10.5), Inches(0.8), [
        "Substituir Telnet por SSH — encripta todo o tráfego",
        "No Wireshark, SSH mostra apenas dados cifrados (ruído)",
    ], bullet_color=ACCENT_GREEN, font_size=15)


def slide_demo_hashes(prs):
    """Slide 10 — Demo Hash Cracking"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Demo: Hash Cracking — CPU vs GPU", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT_ORANGE)

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
                 "15 passwords  ×  4 algoritmos  =  60 hashes  ×  5 modos de ataque  ×  2 dispositivos",
                 font_size=16, color=ACCENT_CYAN, bold=True)

    # Command
    add_shape_fill(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.6), DARK_GRAY)
    add_text_box(slide, Inches(1.2), Inches(2.25), Inches(10.5), Inches(0.5),
                 "$ python orchestrator.py --config config/apresentacao_final.yaml",
                 font_size=15, color=ACCENT_GREEN, font_name="Consolas")

    # GPU results
    add_text_box(slide, Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.4),
                 "Resultados GPU", font_size=18, color=ACCENT_ORANGE, bold=True)

    gpu_output = (
        "[GPU] MD5:     10/15 — 1.8s\n"
        "[GPU] SHA-256:  9/15 — 4.2s\n"
        "[GPU] Bcrypt:   7/15 — 12s\n"
        "[GPU] Argon2:   6/15 — 28s"
    )
    add_shape_fill(slide, Inches(0.8), Inches(3.6), Inches(5.3), Inches(2.2), DARK_GRAY)
    add_text_box(slide, Inches(1.1), Inches(3.7), Inches(5), Inches(2),
                 gpu_output, font_size=15, color=ACCENT_GREEN, font_name="Consolas")

    # CPU results
    add_text_box(slide, Inches(6.8), Inches(3.1), Inches(5.5), Inches(0.4),
                 "Resultados CPU", font_size=18, color=ACCENT_CYAN, bold=True)

    cpu_output = (
        "[CPU] MD5:     10/15 — 12s\n"
        "[CPU] SHA-256:  9/15 — 28s\n"
        "[CPU] Bcrypt:   7/15 — 85s\n"
        "[CPU] Argon2:   6/15 — timeout"
    )
    add_shape_fill(slide, Inches(6.8), Inches(3.6), Inches(5.3), Inches(2.2), DARK_GRAY)
    add_text_box(slide, Inches(7.1), Inches(3.7), Inches(5), Inches(2),
                 cpu_output, font_size=15, color=ACCENT_RED, font_name="Consolas")

    # Bottom note
    add_shape_fill(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.9), RGBColor(0x0A, 0x14, 0x0A))
    add_text_box(slide, Inches(1.2), Inches(6.2), Inches(10.5), Inches(0.7),
                 "As 5 passwords fortes (ex: X7k#mP9$vL2@) resistiram a TODOS os 5 modos\n"
                 "de ataque em TODOS os 4 algoritmos — GPU e CPU.",
                 font_size=14, color=ACCENT_GREEN)


def slide_results(prs):
    """Slide 11 — Resultados Comparativos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Resultados Comparativos", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    data = [
        ["Algoritmo", "Total", "Crackeadas", "Taxa", "Tempo GPU", "Tempo CPU", "Speedup"],
        ["MD5",        "15",   "10",         "66.7%", "1.8s",     "12.0s",     "6.7×"],
        ["SHA-256",    "15",   "9",          "60.0%", "4.2s",     "28.0s",     "6.7×"],
        ["Bcrypt",     "15",   "7",          "46.7%", "12.0s",    "85.0s",     "7.1×"],
        ["Argon2id",   "15",   "6",          "40.0%", "28.0s",    "timeout",   ">6×"],
        ["TOTAL",      "60",   "32",         "53.3%", "~46s",     ">180s",     "—"],
    ]

    add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(3),
              len(data), len(data[0]), data)

    # Key insights
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
                 "Insights Principais", font_size=20, color=WHITE, bold=True)

    insights = [
        "GPU é 6–7× mais rápida que CPU em todos os algoritmos",
        "MD5: 22.5 GH/s — completamente inadequado para passwords",
        "Argon2 neutraliza a vantagem de hardware (memory-hard)",
        "Taxa global de 53% — password fraca = sem proteção independente do algoritmo",
        "5 passwords fortes: 0% crackeadas — segurança = algoritmo forte + password forte",
    ]

    colors = [ACCENT_CYAN, ACCENT_RED, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_GREEN]
    y = Inches(5.5)
    for insight, col in zip(insights, colors):
        add_text_box(slide, Inches(1.0), y, Inches(0.3), Inches(0.3),
                     "▸", font_size=14, color=col)
        add_text_box(slide, Inches(1.4), y, Inches(11), Inches(0.3),
                     insight, font_size=14, color=LIGHT_GRAY)
        y += Inches(0.35)


def slide_security_features(prs):
    """Slide 12 — Segurança Operacional"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Segurança Operacional", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    features = [
        ("Rede Isolada", "Sem rota default — NetworkManager valida\nprogramaticamente antes de cada run",
         ACCENT_CYAN),
        ("Limpeza Segura", "3-pass overwrite de ficheiros sensíveis\n(hashes, passwords, potfiles)",
         ACCENT_RED),
        ("Reprodutibilidade", "Seed determinística + salts fixos\n= mesmos hashes em runs idênticos",
         ACCENT_GREEN),
        ("Separação de Dados", "safe_hashes.py separa hashes de\npasswords em plaintext",
         ACCENT_PURPLE),
        ("Anonimização", "Logs auto-anonimizados\nsem dados pessoais",
         ACCENT_ORANGE),
        ("Config Validada", "Schema YAML validado antes\nde qualquer execução",
         ACCENT_CYAN),
    ]

    x_positions = [Inches(0.5), Inches(4.5), Inches(8.5)]
    y_positions = [Inches(1.8), Inches(4.2)]

    for i, (title, desc, color) in enumerate(features):
        x = x_positions[i % 3]
        y = y_positions[i // 3]
        add_card(slide, x, y, Inches(3.5), Inches(1.7), title,
                 desc.split('\n'), accent=color)


def slide_dev_timeline(prs):
    """Slide 13 — Timeline de Desenvolvimento"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Timeline de Desenvolvimento", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    phases = [
        ("v0.1", "02 Fev", "Fundação", "Pipeline base, WiFi, Telnet, setup scripts", ACCENT_CYAN),
        ("v0.2", "06 Fev", "Hardening", "Windows Defender, validação, documentação", ACCENT_GREEN),
        ("v0.3", "07–08 Fev", "Integração", "Multi-vetor, SHA-256 salted, regras hashcat", ACCENT_ORANGE),
        ("v0.4", "09 Fev", "Apresentação", "15 pw × 5 modos, rockyou 14.3M, bug fixes", ACCENT_PURPLE),
    ]

    y = Inches(1.8)
    # Timeline line
    add_shape_fill(slide, Inches(2.2), Inches(1.8), Pt(3), Inches(5), ACCENT_CYAN)

    for version, date, title, desc, color in phases:
        # Dot on timeline
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.05), y + Inches(0.1), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Date
        add_text_box(slide, Inches(0.5), y + Inches(0.05), Inches(1.5), Inches(0.35),
                     date, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

        # Content
        add_text_box(slide, Inches(2.8), y, Inches(3), Inches(0.35),
                     f"{version} — {title}", font_size=18, color=color, bold=True)
        add_text_box(slide, Inches(2.8), y + Inches(0.4), Inches(9), Inches(0.4),
                     desc, font_size=14, color=LIGHT_GRAY)

        y += Inches(1.3)

    # Stats
    add_text_box(slide, Inches(7.5), Inches(2.0), Inches(5), Inches(0.4),
                 "Estatísticas", font_size=20, color=WHITE, bold=True)

    stats = [
        ("35", "commits"),
        ("8", "dias de desenvolvimento"),
        ("12", "ficheiros Python"),
        ("~4500", "linhas de código"),
        ("7", "algoritmos de hashing"),
        ("4", "modos de ataque"),
        ("3", "SO suportados"),
    ]

    y = Inches(2.6)
    for val, label in stats:
        add_text_box(slide, Inches(7.5), y, Inches(1), Inches(0.3),
                     val, font_size=16, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, Inches(8.5), y, Inches(4), Inches(0.3),
                     label, font_size=14, color=LIGHT_GRAY)
        y += Inches(0.35)


def slide_conclusions(prs):
    """Slide 14 — Conclusões"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Conclusões", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3))

    conclusions = [
        ("WiFi WPA2", "Password crackeada em 3 segundos.\nProteção: passwords longas e aleatórias (≥16 chars).",
         ACCENT_GREEN, "3s"),
        ("Telnet", "Credenciais visíveis em texto claro no Wireshark.\nProteção: usar SSH em vez de Telnet.",
         ACCENT_PURPLE, "0 encriptação"),
        ("MD5", "22.5 GH/s na GPU — completamente inadequado.\n16× mais rápido que CPU.",
         ACCENT_RED, "22.5 GH/s"),
        ("Argon2", "Memory-hard — neutraliza vantagem da GPU.\nTimeout no CPU prova eficácia.",
         ACCENT_CYAN, "140 H/s"),
    ]

    x = Inches(0.5)
    for title, desc, color, metric in conclusions:
        add_shape_fill(slide, x, Inches(1.8), Inches(2.9), Inches(2.8), DARK_GRAY)
        add_shape_fill(slide, x, Inches(1.8), Inches(2.9), Pt(3), color)
        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.3),
                     title, font_size=16, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.4),
                     metric, font_size=28, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(1.2),
                     desc, font_size=11, color=LIGHT_GRAY)
        x += Inches(3.1)

    # Final message
    add_shape_fill(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.2), DARK_GRAY)
    add_shape_fill(slide, Inches(1.5), Inches(5.2), Inches(10.3), Pt(3), ACCENT_CYAN)
    add_text_box(slide, Inches(1.5), Inches(5.45), Inches(10.3), Inches(0.5),
                 "Segurança = Algoritmo Forte + Password Forte",
                 font_size=28, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.95), Inches(10.3), Inches(0.4),
                 "Sem os dois em simultâneo, estamos vulneráveis.",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_thanks(prs):
    """Slide 15 — Obrigado + Perguntas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_CYAN)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
                 "Obrigado!", font_size=52, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.5),
                 "Perguntas?", font_size=30, color=ACCENT_CYAN,
                 alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.0), Inches(2.3))

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
                 "Henrique Carvalho  ·  Gonçalo Ferro  ·  Francisco Silva  ·  Duarte Vilar",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
                 "Projeto Final — Cibersegurança  |  Fevereiro 2026",
                 font_size=14, color=RGBColor(0x70, 0x78, 0x80),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
                 "github.com/Henryu1781/HashCrackerLab",
                 font_size=14, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    add_shape_fill(slide, Inches(0), Inches(7.2), SLIDE_W, Pt(4), ACCENT_CYAN)


# ═══════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)           # 1
    slide_agenda(prs)          # 2
    slide_context(prs)         # 3
    slide_architecture(prs)    # 4
    slide_pipeline(prs)        # 5
    slide_algorithms(prs)      # 6
    slide_attack_modes(prs)    # 7
    slide_demo_wifi(prs)       # 8
    slide_demo_telnet(prs)     # 9
    slide_demo_hashes(prs)     # 10
    slide_results(prs)         # 11
    slide_security_features(prs)  # 12
    slide_dev_timeline(prs)    # 13
    slide_conclusions(prs)     # 14
    slide_thanks(prs)          # 15

    output = "HashCrackerLab_Apresentacao.pptx"
    prs.save(output)
    print(f"[+] Apresentação gerada: {output}")
    print(f"    → 15 slides")
    print(f"    → Tema: dark + ciano")
    print(f"    → Widescreen 16:9")


if __name__ == "__main__":
    main()
